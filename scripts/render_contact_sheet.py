#!/usr/bin/env python3
"""Render PPTX slide thumbnails and a contact sheet.

Backends:
- quicklook: macOS QuickLook, no Office installation required.
- libreoffice: LibreOffice -> PDF -> PyMuPDF render, works on macOS/Linux/Windows.
- powerpoint: Microsoft PowerPoint COM export, Windows only.
"""
from __future__ import annotations

import argparse
import platform
import shutil
import subprocess
from pathlib import Path

from PIL import Image, ImageDraw, ImageOps
from pptx import Presentation


def split_slides(src: Path, out_dir: Path) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    base = Presentation(str(src))
    count = len(base.slides)
    for keep in range(count):
        prs = Presentation(str(src))
        slide_ids = prs.slides._sldIdLst  # noqa: SLF001
        for idx in range(count - 1, -1, -1):
            if idx == keep:
                continue
            slide_id = slide_ids[idx]
            prs.part.drop_rel(slide_id.rId)
            slide_ids.remove(slide_id)
        prs.save(out_dir / f"slide_{keep + 1:02d}.pptx")
    return count


def render_quicklook(src: Path, out_dir: Path, size: int) -> int:
    single_dir = out_dir / "single_slides"
    preview_dir = out_dir / "previews"
    count = split_slides(src, single_dir)
    preview_dir.mkdir(parents=True, exist_ok=True)
    qlmanage = shutil.which("qlmanage")
    if not qlmanage:
        raise RuntimeError("qlmanage not found.")
    for pptx in sorted(single_dir.glob("slide_*.pptx")):
        subprocess.run(
            [qlmanage, "-t", "-s", str(size), "-o", str(preview_dir), str(pptx)],
            check=False,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    rendered = 0
    for png in sorted(preview_dir.glob("slide_*.pptx.png")):
        slide_no = png.stem.replace(".pptx", "").split("_")[-1]
        png.rename(preview_dir / f"slide_{slide_no}.png")
        rendered += 1
    return rendered


def _find_soffice() -> str | None:
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    candidates = [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return candidate
    return None


def render_libreoffice(src: Path, out_dir: Path, dpi: int) -> int:
    soffice = _find_soffice()
    if not soffice:
        raise RuntimeError("LibreOffice executable not found.")
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required for the libreoffice backend: pip install pymupdf") from exc

    pdf_dir = out_dir / "pdf"
    preview_dir = out_dir / "previews"
    pdf_dir.mkdir(parents=True, exist_ok=True)
    preview_dir.mkdir(parents=True, exist_ok=True)

    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(pdf_dir), str(src)],
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
    )
    pdf_path = pdf_dir / f"{src.stem}.pdf"
    if not pdf_path.exists():
        candidates = sorted(pdf_dir.glob("*.pdf"))
        if not candidates:
            raise RuntimeError("LibreOffice did not produce a PDF.")
        pdf_path = candidates[0]

    doc = fitz.open(pdf_path)
    zoom = dpi / 72
    matrix = fitz.Matrix(zoom, zoom)
    for idx, page in enumerate(doc, 1):
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        pix.save(preview_dir / f"slide_{idx:02d}.png")
    return len(doc)


def render_powerpoint(src: Path, out_dir: Path, width: int) -> int:
    if platform.system() != "Windows":
        raise RuntimeError("PowerPoint backend is Windows-only.")
    try:
        import win32com.client  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError("PowerPoint backend requires pywin32: pip install pywin32") from exc

    preview_dir = out_dir / "previews"
    preview_dir.mkdir(parents=True, exist_ok=True)
    app = win32com.client.Dispatch("PowerPoint.Application")
    presentation = app.Presentations.Open(str(src.resolve()), WithWindow=False)
    try:
        presentation.Export(str(preview_dir.resolve()), "PNG", width)
        exported = sorted(preview_dir.glob("Slide*.PNG"))
        for idx, path in enumerate(exported, 1):
            path.rename(preview_dir / f"slide_{idx:02d}.png")
        return len(exported)
    finally:
        presentation.Close()
        app.Quit()


def choose_backend(requested: str) -> str:
    if requested != "auto":
        return requested
    if platform.system() == "Darwin" and shutil.which("qlmanage"):
        return "quicklook"
    if _find_soffice():
        return "libreoffice"
    if platform.system() == "Windows":
        return "powerpoint"
    raise RuntimeError("No render backend found. Install LibreOffice or use macOS QuickLook/Windows PowerPoint.")


def make_contact_sheet(preview_dir: Path, out_path: Path, thumb_width: int = 300, cols: int = 4) -> int:
    paths = sorted(preview_dir.glob("slide_*.png"))
    thumbs = []
    for path in paths:
        image = Image.open(path).convert("RGB")
        image = ImageOps.expand(image, border=8, fill="white")
        draw = ImageDraw.Draw(image)
        label = path.stem.split("_")[-1]
        draw.rectangle((8, 8, 72, 38), fill=(255, 255, 255), outline=(60, 90, 150))
        draw.text((18, 16), label, fill=(40, 70, 130))
        ratio = thumb_width / image.width
        thumbs.append(image.resize((thumb_width, int(image.height * ratio))))

    if not thumbs:
        raise RuntimeError("No preview PNGs generated.")

    rows = (len(thumbs) + cols - 1) // cols
    row_height = max(thumb.height for thumb in thumbs)
    sheet = Image.new("RGB", (cols * thumb_width, rows * row_height), (238, 242, 247))
    for idx, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((idx % cols) * thumb_width, (idx // cols) * row_height))
    sheet.save(out_path, quality=92)
    return len(thumbs)


def main() -> int:
    parser = argparse.ArgumentParser(description="Render PPTX slide thumbnails and a contact sheet.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--out-dir", type=Path, default=Path("tmp/ppt_preview"))
    parser.add_argument("--backend", choices=("auto", "quicklook", "libreoffice", "powerpoint"), default="auto")
    parser.add_argument("--size", type=int, default=1400, help="QuickLook/PowerPoint export width in pixels")
    parser.add_argument("--dpi", type=int, default=160, help="LibreOffice PDF render DPI")
    parser.add_argument("--thumb-width", type=int, default=300)
    parser.add_argument("--cols", type=int, default=4)
    args = parser.parse_args()

    if args.out_dir.exists():
        shutil.rmtree(args.out_dir)
    args.out_dir.mkdir(parents=True, exist_ok=True)

    backend = choose_backend(args.backend)
    if backend == "quicklook":
        rendered = render_quicklook(args.pptx, args.out_dir, args.size)
    elif backend == "libreoffice":
        rendered = render_libreoffice(args.pptx, args.out_dir, args.dpi)
    else:
        rendered = render_powerpoint(args.pptx, args.out_dir, args.size)

    contact_sheet = args.out_dir / "contact_sheet.jpg"
    sheet_count = make_contact_sheet(args.out_dir / "previews", contact_sheet, args.thumb_width, args.cols)

    print(f"backend: {backend}")
    print(f"rendered: {rendered}")
    print(f"contact_sheet_slides: {sheet_count}")
    print(f"contact_sheet: {contact_sheet}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
