#!/usr/bin/env python3
"""Analyze a PPTX template and summarize reusable visual style signals."""
from __future__ import annotations

import argparse
import json
import statistics
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation

EMU_PER_INCH = 914400


def inches(v: int) -> float:
    return round(v / EMU_PER_INCH, 3)


def color_to_hex(color: Any) -> str | None:
    if color is None:
        return None
    try:
        rgb = color.rgb
    except (AttributeError, ValueError):
        return None
    if rgb is None:
        return None
    return f"#{rgb}"


def shape_text(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text_frame.text.strip()
    return ""


def add_run_stats(shape: Any, fonts: Counter[str], font_sizes: Counter[float], colors: Counter[str]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            if run.font.name:
                fonts[run.font.name] += 1
            if run.font.size:
                font_sizes[round(run.font.size.pt, 1)] += 1
            color = color_to_hex(run.font.color)
            if color:
                colors[color] += 1


def add_shape_colors(shape: Any, colors: Counter[str]) -> None:
    for attr in ("fill", "line"):
        obj = getattr(shape, attr, None)
        if obj is None:
            continue
        try:
            color = color_to_hex(obj.fore_color)
        except (AttributeError, TypeError, ValueError):
            color = None
        if color:
            colors[color] += 1


def classify_position(x: float, y: float, w: float, h: float, slide_w: float, slide_h: float) -> list[str]:
    tags = []
    cx = x + w / 2
    cy = y + h / 2
    if y < slide_h * 0.18:
        tags.append("top/title-area")
    if cy > slide_h * 0.82:
        tags.append("bottom/footer-area")
    if cx < slide_w * 0.18:
        tags.append("left-edge")
    if cx > slide_w * 0.82:
        tags.append("right-edge")
    if abs(cx - slide_w / 2) < slide_w * 0.08:
        tags.append("center-axis")
    return tags


def summarize_positions(values: list[float]) -> dict[str, float] | None:
    if not values:
        return None
    return {
        "min": round(min(values), 3),
        "median": round(statistics.median(values), 3),
        "max": round(max(values), 3),
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Analyze a PPTX template's reusable visual style.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--json-out", type=Path)
    parser.add_argument("--top", type=int, default=12)
    args = parser.parse_args()

    prs = Presentation(str(args.pptx))
    slide_w = inches(prs.slide_width)
    slide_h = inches(prs.slide_height)

    fonts: Counter[str] = Counter()
    font_sizes: Counter[float] = Counter()
    colors: Counter[str] = Counter()
    picture_positions = []
    text_positions = []
    title_candidates = []
    footer_candidates = []
    slide_summaries = []
    shape_type_counter: Counter[str] = Counter()
    masters = len(prs.slide_masters)
    layouts = len(prs.slide_layouts)

    for slide_no, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide": slide_no,
            "layout": slide.slide_layout.name,
            "shapes": len(slide.shapes),
            "text_preview": [],
            "pictures": 0,
        }
        for shape_no, shape in enumerate(slide.shapes, 1):
            x, y, w, h = map(inches, (shape.left, shape.top, shape.width, shape.height))
            shape_type = str(shape.shape_type)
            shape_type_counter[shape_type] += 1
            add_run_stats(shape, fonts, font_sizes, colors)
            add_shape_colors(shape, colors)
            text = shape_text(shape)
            tags = classify_position(x, y, w, h, slide_w, slide_h)

            if text:
                text_positions.append({"slide": slide_no, "shape": shape_no, "x": x, "y": y, "w": w, "h": h, "tags": tags, "text": text[:80]})
                if "top/title-area" in tags and h < slide_h * 0.18:
                    title_candidates.append(text_positions[-1])
                if "bottom/footer-area" in tags and h < slide_h * 0.12:
                    footer_candidates.append(text_positions[-1])
                if len(slide_info["text_preview"]) < 5:
                    slide_info["text_preview"].append(text[:60].replace("\n", " / "))

            if shape.shape_type == 13:
                slide_info["pictures"] += 1
                picture_positions.append({"slide": slide_no, "shape": shape_no, "x": x, "y": y, "w": w, "h": h, "tags": tags})

        slide_summaries.append(slide_info)

    title_y = summarize_positions([item["y"] for item in title_candidates])
    footer_y = summarize_positions([item["y"] for item in footer_candidates])
    picture_x = summarize_positions([item["x"] for item in picture_positions])
    picture_y = summarize_positions([item["y"] for item in picture_positions])

    report = {
        "file": str(args.pptx),
        "slide_count": len(prs.slides),
        "slide_size_in": {"width": slide_w, "height": slide_h},
        "masters": masters,
        "layouts": layouts,
        "top_fonts": fonts.most_common(args.top),
        "top_font_sizes_pt": font_sizes.most_common(args.top),
        "top_colors": colors.most_common(args.top),
        "shape_types": shape_type_counter.most_common(),
        "position_hints": {
            "title_y_in": title_y,
            "footer_y_in": footer_y,
            "picture_x_in": picture_x,
            "picture_y_in": picture_y,
        },
        "title_candidates": title_candidates[: args.top],
        "footer_candidates": footer_candidates[: args.top],
        "picture_candidates": picture_positions[: args.top],
        "slides": slide_summaries,
    }

    if args.json_out:
        args.json_out.parent.mkdir(parents=True, exist_ok=True)
        args.json_out.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"file: {report['file']}")
    print(f"slides: {report['slide_count']}")
    print(f"size: {slide_w} x {slide_h} in")
    print(f"masters/layouts: {masters}/{layouts}")
    print(f"top_fonts: {report['top_fonts']}")
    print(f"top_font_sizes_pt: {report['top_font_sizes_pt']}")
    print(f"top_colors: {report['top_colors']}")
    print(f"position_hints: {report['position_hints']}")
    if args.json_out:
        print(f"json: {args.json_out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
