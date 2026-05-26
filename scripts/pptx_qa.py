#!/usr/bin/env python3
"""Structural QA for generated defense PPTX decks."""
from __future__ import annotations

import argparse
import statistics
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation

EMU_PER_INCH = 914400


def inches(v: int) -> float:
    return v / EMU_PER_INCH


def shape_text(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text_frame.text.strip().replace("\n", " / ")
    return ""


def is_picture(shape: Any) -> bool:
    return shape.shape_type == 13


def is_background_like(shape: Any, slide_w: int, slide_h: int) -> bool:
    return shape.width > slide_w * 0.95 and shape.height > slide_h * 0.95


def shape_area(shape: Any) -> int:
    return int(shape.width) * int(shape.height)


def non_empty_text_length(shape: Any) -> int:
    return len(shape_text(shape).replace(" ", "").replace("/", ""))


def collect_font_sizes(shape: Any) -> list[float]:
    sizes = []
    if not getattr(shape, "has_text_frame", False):
        return sizes
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip() and run.font.size:
                sizes.append(round(run.font.size.pt, 1))
    return sizes


def find_alignment_groups(edges: list[tuple[int, int, float]], tolerance: float) -> list[tuple[float, list[tuple[int, int]]]]:
    groups: list[list[tuple[int, int, float]]] = []
    for item in sorted(edges, key=lambda x: x[2]):
        for group in groups:
            if abs(group[0][2] - item[2]) <= tolerance:
                group.append(item)
                break
        else:
            groups.append([item])
    result = []
    for group in groups:
        if len(group) >= 3:
            result.append((round(statistics.median([item[2] for item in group]), 3), [(item[0], item[1]) for item in group]))
    return result


def main() -> int:
    parser = argparse.ArgumentParser(description="Check common PPTX layout issues.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--expected-slides", type=int, default=None)
    parser.add_argument("--tiny-width", type=float, default=0.13, help="tiny text width threshold in inches")
    parser.add_argument("--tiny-height", type=float, default=0.08, help="tiny text height threshold in inches")
    parser.add_argument("--edge-tolerance", type=float, default=0.035, help="alignment edge grouping tolerance in inches")
    parser.add_argument("--center-tolerance", type=float, default=0.08, help="center-axis warning tolerance in inches")
    parser.add_argument("--max-font-spread", type=float, default=9.5, help="warn if one slide uses too many font sizes")
    parser.add_argument("--warn-only", action="store_true", help="print warnings but exit 0 unless hard errors occur")
    args = parser.parse_args()

    prs = Presentation(str(args.pptx))
    width, height = prs.slide_width, prs.slide_height
    center_x = width / 2

    out_of_bounds = []
    tiny_text = []
    bullet_text = []
    overlaps = []
    sparse_large_text_boxes = []
    center_axis_warnings = []
    weak_alignment_pages = []
    font_spread_warnings = []
    font_counter: Counter[float] = Counter()
    page_font_counter: dict[int, Counter[float]] = defaultdict(Counter)

    for slide_no, slide in enumerate(prs.slides, 1):
        pictures = []
        text_boxes = []
        layout_shapes = []
        left_edges = []
        right_edges = []
        font_sizes_on_slide = []

        for shape_no, shape in enumerate(slide.shapes, 1):
            x, y, w, h = shape.left, shape.top, shape.width, shape.height
            text = shape_text(shape)
            area_ratio = shape_area(shape) / (width * height)

            if x < -20000 or y < -20000 or x + w > width + 20000 or y + h > height + 20000:
                out_of_bounds.append((slide_no, shape_no, inches(x), inches(y), inches(w), inches(h), text[:80]))

            if text:
                if "•" in text:
                    bullet_text.append((slide_no, shape_no, text[:120]))
                if w < args.tiny_width * EMU_PER_INCH or h < args.tiny_height * EMU_PER_INCH:
                    tiny_text.append((slide_no, shape_no, inches(w), inches(h), text[:80]))
                if w > 0.25 * EMU_PER_INCH and h > 0.10 * EMU_PER_INCH:
                    text_boxes.append((shape_no, x, y, w, h, text[:40]))
                    layout_shapes.append((shape_no, x, y, w, h, text[:40]))
                    left_edges.append((slide_no, shape_no, inches(x)))
                    right_edges.append((slide_no, shape_no, inches(x + w)))
                text_len = non_empty_text_length(shape)
                if area_ratio > 0.105 and text_len < 45 and y > height * 0.13 and y + h < height * 0.9:
                    sparse_large_text_boxes.append((slide_no, shape_no, round(area_ratio, 3), text_len, text[:90]))

                sizes = collect_font_sizes(shape)
                font_sizes_on_slide.extend(sizes)
                for size in sizes:
                    font_counter[size] += 1
                    page_font_counter[slide_no][size] += 1

            if is_picture(shape):
                if not is_background_like(shape, width, height):
                    pictures.append((shape_no, x, y, w, h))
                    layout_shapes.append((shape_no, x, y, w, h, "[picture]"))
                    left_edges.append((slide_no, shape_no, inches(x)))
                    right_edges.append((slide_no, shape_no, inches(x + w)))

        for text_item in text_boxes:
            text_no, tx, ty, tw, th, preview = text_item
            for pic_no, px, py, pw, ph in pictures:
                ix = max(0, min(tx + tw, px + pw) - max(tx, px))
                iy = max(0, min(ty + th, py + ph) - max(ty, py))
                if ix * iy > min(tw * th, pw * ph) * 0.18:
                    overlaps.append((slide_no, text_no, pic_no, preview))

        if len(layout_shapes) >= 4:
            left_groups = find_alignment_groups(left_edges, args.edge_tolerance)
            right_groups = find_alignment_groups(right_edges, args.edge_tolerance)
            if not left_groups and not right_groups:
                weak_alignment_pages.append((slide_no, len(layout_shapes), "no common left/right edge group of >=3 shapes"))

        if slide_no in {1, len(prs.slides)}:
            centered = []
            for shape_no, x, y, w, h, preview in layout_shapes:
                if y < height * 0.92 and w < width * 0.9:
                    delta = abs((x + w / 2) - center_x)
                    if delta > args.center_tolerance * EMU_PER_INCH:
                        centered.append((shape_no, round(inches(delta), 3), preview))
            if centered:
                center_axis_warnings.append((slide_no, centered[:8]))

        if font_sizes_on_slide:
            spread = max(font_sizes_on_slide) - min(font_sizes_on_slide)
            unique_sizes = sorted(set(font_sizes_on_slide))
            if spread > args.max_font_spread and len(unique_sizes) > 5:
                font_spread_warnings.append((slide_no, round(spread, 1), unique_sizes))

    print(f"slides: {len(prs.slides)}")
    print(f"size: {inches(width):.3f} x {inches(height):.3f} in")
    if args.expected_slides is not None:
        print(f"expected_slides: {args.expected_slides}")
        if len(prs.slides) != args.expected_slides:
            print("ERROR: slide count mismatch")

    print(f"out_of_bounds: {len(out_of_bounds)}")
    for item in out_of_bounds[:20]:
        print("  ", item)

    print(f"tiny_text_boxes: {len(tiny_text)}")
    for item in tiny_text[:20]:
        print("  ", item)

    print(f"bullet_texts: {len(bullet_text)}")
    for item in bullet_text[:20]:
        print("  ", item)

    print(f"rough_text_picture_overlaps: {len(overlaps)}")
    for item in overlaps[:20]:
        print("  ", item)

    print(f"sparse_large_text_boxes: {len(sparse_large_text_boxes)}")
    for item in sparse_large_text_boxes[:20]:
        print("  ", item)

    print(f"center_axis_warnings: {len(center_axis_warnings)}")
    for item in center_axis_warnings[:10]:
        print("  ", item)

    print(f"weak_alignment_pages: {len(weak_alignment_pages)}")
    for item in weak_alignment_pages[:20]:
        print("  ", item)

    print(f"font_spread_warnings: {len(font_spread_warnings)}")
    for item in font_spread_warnings[:20]:
        print("  ", item)

    print("font_sizes:", sorted(font_counter.items()))
    print("page_font_sizes:")
    for slide_no in sorted(page_font_counter):
        print(f"  P{slide_no:02d}: {sorted(page_font_counter[slide_no].items())}")

    hard_failed = False
    if args.expected_slides is not None and len(prs.slides) != args.expected_slides:
        hard_failed = True
    if out_of_bounds or tiny_text or bullet_text or overlaps:
        hard_failed = True

    warn_failed = bool(sparse_large_text_boxes or center_axis_warnings or weak_alignment_pages or font_spread_warnings)
    return 1 if hard_failed or (warn_failed and not args.warn_only) else 0


if __name__ == "__main__":
    raise SystemExit(main())
